#!/usr/bin/env python3
"""生成答辩 PPTX — 诚实反映 4 层沙箱隔离架构.

安装依赖（已安装则跳过）:
    pip install python-pptx

用法:
    python scripts/_gen_ppt.py

输出:
    - ~/Desktop/security-agent-答辩.pptx
    - docs/competitions/security-agent-答辩.pptx
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ── 团队信息 ─────────────────────────────────
TEAM_NAME = "智能麒麟agent安全运维"
LEADER = "欧阳学武"
MEMBER = "黄照庭"
DATE_RANGE = "2026.4.17 - 2026.7.19"

# ── 配色 ──────────────────────────────────
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)  # 深蓝黑背景
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)  # 强调蓝
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)  # 强调绿
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)

# 4 层沙箱描述（与 profile.py 实际一致）
SANDBOX_4_LAYER_DESC = (
    "4层沙箱隔离: setuid降权 + rlimit资源限制 + "
    "OverlayFS写时复制 + mount_ns文件隔离"
)

# ── 辅助函数 ───────────────────────────────


def set_slide_bg(slide, color: RGBColor):
    """设置幻灯片纯色背景."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide, left: float, top: float, width: float, height: float,
    text: str, font_size: int = 18, bold: bool = False,
    color: RGBColor = WHITE, alignment: int = PP_ALIGN.LEFT,
    font_name: str = "微软雅黑",
) -> None:
    """在幻灯片上添加文本框."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(
    slide, left: float, top: float, width: float, height: float,
    lines: list[tuple[str, int, bool, RGBColor]],
    alignment: int = PP_ALIGN.LEFT,
) -> None:
    """多行文本框，每行指定 (text, font_size, bold, color)."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.alignment = alignment
        p.space_after = Pt(6)


# ── 幻灯片生成 ─────────────────────────────


def create_cover(slide):
    """P1 封面."""
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.8, 1.2, 8.4, 1.2, "银河麒麟智能安全运维 Agent",
                font_size=36, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, 2.4, 8.4, 0.6, "security-agent v0.9.0",
                font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, 3.6, 8.4, 0.5, f"队名: {TEAM_NAME}",
                font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, 4.1, 8.4, 0.5, f"队长: {LEADER}  |  队员: {MEMBER}",
                font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, 4.6, 8.4, 0.5, f"日期: {DATE_RANGE}",
                font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, 5.8, 8.4, 0.5,
                "第十五届\"中国软件杯\"大学生软件设计大赛 · A2赛题 · 麒麟软件有限公司",
                font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def create_toc(slide):
    """P2 目录."""
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "目录",
                font_size=28, bold=True, color=ACCENT_BLUE)
    items = [
        "1. 问题背景",
        "2. 产品定位",
        "3. 技术架构 — 五层流水线 + 三 Agent",
        "4. MCP 插件化",
        "5. 三层安全防御",
        "6. 4 层沙箱隔离 ← 核心",
        "7. 人工审批闭环",
        "8. 智能终端",
        "9. 全链路溯源",
        "10. 麒麟 LoongArch 适配",
        "11. 测试与质量保障",
    ]
    for i, item in enumerate(items):
        color = YELLOW if "4 层沙箱" in item else WHITE
        add_textbox(slide, 1.0, 1.2 + i * 0.45, 8, 0.4, item,
                    font_size=16, color=color)


def create_tech_arch(slide):
    """P3 技术架构."""
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "技术架构",
                font_size=26, bold=True, color=ACCENT_BLUE)
    lines = [
        ("┌─ L1 多维感知 ─────────────────────────┐", 14, False, WHITE),
        ("│  8维仪表盘 · 边界对抗 · 知识检索      │", 14, False, WHITE),
        ("├─ L2 安全防护 ─────────────────────────┤", 14, False, ACCENT_GREEN),
        ("│  三层防御 30/35/35 · 4层沙箱          │", 14, True, ACCENT_GREEN),
        ("├────── GATE 层间门禁 ──────────────────┤", 14, False, WHITE),
        ("├─ L3 工具执行 ─────────────────────────┤", 14, False, WHITE),
        ("│  17 Skills · 四工具簇 · 智能终端      │", 14, False, WHITE),
        ("├─ L4 审计追溯 ─────────────────────────┤", 14, False, WHITE),
        ("│  IncidentSpine · append-only 卷宗     │", 14, False, WHITE),
        ("├─ L5 量化迭代 ─────────────────────────┤", 14, False, WHITE),
        ("│  六维指标 · 策略反写 L1               │", 14, False, WHITE),
        ("└──────────────────────────────────────┘", 14, False, WHITE),
    ]
    add_multiline_textbox(slide, 1.5, 1.2, 7, 4.5, lines)
    add_textbox(slide, 0.5, 6.0, 9, 0.4,
                "三 Agent: core_dispatch(L1+L3) + safety_sandbox(L2) + audit_iteration(L4+L5)",
                font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def create_three_layer_defense(slide):
    """P4 三层安全防御."""
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "三层安全防御 — 让 AI 不能乱来",
                font_size=24, bold=True, color=ACCENT_BLUE)
    lines = [
        ("用户输入命令", 16, True, WHITE),
        ("    ↓", 14, False, LIGHT_GRAY),
        ("┌─ 第1层：静态风险评估 (30%) ────────┐", 13, False, WHITE),
        ("│  规则引擎 · 注入扫描 · 四级风险矩阵  │", 13, False, WHITE),
        ("└────────────────────────────────────┘", 13, False, WHITE),
        ("    ↓", 14, False, LIGHT_GRAY),
        ("┌─ 第2层：动态意图审计 (35%) ────────┐", 13, False, WHITE),
        ("│  用户意图 vs 拟执行命令 · 交叉校验   │", 13, False, WHITE),
        ("└────────────────────────────────────┘", 13, False, WHITE),
        ("    ↓", 14, False, LIGHT_GRAY),
        ("┌─ 第3层：受限执行环境 (35%) ────────┐", 13, False, ACCENT_GREEN),
        ("│  4层沙箱 · OverlayFS COW · 最小权限  │", 13, True, ACCENT_GREEN),
        ("└────────────────────────────────────┘", 13, False, ACCENT_GREEN),
        ("    ↓", 14, False, LIGHT_GRAY),
        ("裁定: ALLOW / CONFIRM / APPROVE / ESCALATE / DENY", 14, True, YELLOW),
    ]
    add_multiline_textbox(slide, 1.0, 1.1, 8, 5.5, lines)


def create_sandbox_4layer(slide):
    """P5 4层沙箱隔离 — 核心页."""
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "4 层沙箱隔离 — 写得进去，回得来",
                font_size=24, bold=True, color=ACCENT_BLUE)

    # 四层递进
    layers = [
        ("Layer 0", "setuid 降权", "最低权限用户执行，杜绝提权", ACCENT_GREEN),
        ("Layer 1", "rlimit 资源限制", "CPU/内存/进程数硬限制，防止资源耗尽", ACCENT_GREEN),
        ("Layer 2", "OverlayFS 写时复制", "所有写操作落入 upperdir，原始文件不受影响", ACCENT_GREEN),
        ("Layer 3", "mount_ns 文件隔离", "私有 /tmp、/dev、/proc 挂载点，隔离文件系统", ACCENT_GREEN),
    ]
    for i, (layer, title, desc, color) in enumerate(layers):
        y = 1.2 + i * 1.3
        add_textbox(slide, 1.0, y, 2.0, 0.4, layer,
                    font_size=16, bold=True, color=color)
        add_textbox(slide, 3.0, y, 2.5, 0.4, title,
                    font_size=16, bold=True, color=WHITE)
        add_textbox(slide, 3.0, y + 0.4, 5.5, 0.6, desc,
                    font_size=12, color=LIGHT_GRAY)

    # 底部说明
    add_textbox(slide, 0.5, 6.5, 9, 0.4,
                "注意: network_ns、seccomp、cgroup 为 P1 规划，当前未启用",
                font_size=11, color=YELLOW)


def create_innovation(slide):
    """P6 创新点."""
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 0.3, 9, 0.6, "6 项架构创新",
                font_size=26, bold=True, color=ACCENT_BLUE)
    items = [
        ("4 层沙箱", "setuid降权 + rlimit + OverlayFS COW + mount_ns文件隔离"),
        ("能力装箱", "ToolBox/FlowBox/PluginBox 统一入口 + 熔断保护"),
        ("边界 Fuzzer", "12 探针 + 7 变异策略，自动检测沙箱穿透"),
        ("三方统一契约", "前端/后端/文档共享一份 JSON 真源"),
        ("文档活化", "TF-IDF 双索引（BM25+余弦），零外部依赖"),
        ("知识自愈", "一致性检查 + 新鲜度检测 + 休眠告警"),
    ]
    for i, (title, desc) in enumerate(items):
        y = 1.2 + i * 0.85
        is_sandbox = "4 层沙箱" in title
        add_textbox(slide, 1.0, y, 2.5, 0.4, title,
                    font_size=16, bold=True,
                    color=ACCENT_GREEN if is_sandbox else WHITE)
        add_textbox(slide, 3.5, y, 5.5, 0.4, desc,
                    font_size=13, color=LIGHT_GRAY)


def create_cover_page(slide, title: str, subtitle: str = ""):
    """通用章节封面."""
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 2.5, 9, 1.0, title,
                font_size=30, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, 0.5, 3.6, 9, 0.6, subtitle,
                    font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def create_closing(slide):
    """尾页."""
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, 0.5, 2.0, 9, 0.8, "感谢聆听",
                font_size=36, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 3.0, 9, 0.5, f"队名: {TEAM_NAME}",
                font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 3.5, 9, 0.5, f"队长: {LEADER}  |  队员: {MEMBER}",
                font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 4.5, 9, 0.5,
                "GitHub: github.com/23oyxw/security-agent",
                font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 5.5, 9, 0.5,
                "让安全运维更智能",
                font_size=16, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)


# ── 主函数 ─────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 获取空白布局
    blank_layout = prs.slide_layouts[6]  # Blank

    # ── 按顺序生成幻灯片 ──
    slides_builders = [
        ("cover", lambda: create_cover(prs.slides.add_slide(blank_layout))),
        ("toc", lambda: create_toc(prs.slides.add_slide(blank_layout))),
        ("tech_arch", lambda: create_tech_arch(prs.slides.add_slide(blank_layout))),
        ("three_layer", lambda: create_three_layer_defense(prs.slides.add_slide(blank_layout))),
        ("sandbox", lambda: create_sandbox_4layer(prs.slides.add_slide(blank_layout))),
        ("innovation", lambda: create_innovation(prs.slides.add_slide(blank_layout))),
        ("closing", lambda: create_closing(prs.slides.add_slide(blank_layout))),
    ]
    for name, fn in slides_builders:
        fn()
        print(f"  [OK] Slide: {name}")

    # ── 保存 ──
    fname = "security-agent-答辩.pptx"

    # 输出 1: Desktop
    desktop = Path(os.environ.get("USERPROFILE", "~/Desktop")) / fname
    prs.save(str(desktop))
    print(f"\n  => Saved to: {desktop}")

    # 输出 2: docs/competitions/
    repo_docs = Path(__file__).resolve().parent.parent / "docs" / "competitions" / fname
    prs.save(str(repo_docs))
    print(f"  => Saved to: {repo_docs}")

    print(f"\n  总页数: {len(prs.slides)}")
    print(f"  沙箱描述: {SANDBOX_4_LAYER_DESC}")
    print("  Done.")


if __name__ == "__main__":
    main()
